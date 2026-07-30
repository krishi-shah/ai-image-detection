"""Build reports/presentation.pptx for the EECS 4080 final presentation.

Usage:
    python scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "presentation.pptx"
PLOTS = ROOT / "outputs" / "plots"


def _run(paragraph, text: str, size: int = 18, bold: bool = False, color=(30, 30, 30)):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"
    return run


def add_bullet_slide(prs, title: str, bullets: list[str], image: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    _run(title_box.text_frame.paragraphs[0], title, size=26, bold=True, color=(20, 40, 80))

    width = Inches(6.2) if image and image.exists() else Inches(12)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), width, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, f"• {bullet}", size=16)
        p.space_after = Pt(8)

    if image and image.exists():
        slide.shapes.add_picture(str(image), Inches(7.0), Inches(1.3), width=Inches(5.5))
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5))
    tf = tb.text_frame
    _run(tf.paragraphs[0], "AI-Generated Image Detection", 34, True, (20, 40, 80))
    p = tf.add_paragraph()
    _run(p, "Evaluating Generalisation Across Generative Model Generations", 18, False, (60, 60, 60))
    p = tf.add_paragraph()
    _run(p, "Krishi Rajeshkumar Shah (220968905)  ·  Supervisor: Mona Nasery", 14, False, (90, 90, 90))
    p = tf.add_paragraph()
    _run(p, "EECS 4080 · York University · Summer 2026", 14, False, (90, 90, 90))

    add_bullet_slide(prs, "The Problem", [
        "Detectors look strong on older benchmarks (CIFAKE, ProGAN-era sets).",
        "Newer generators (SD3, Midjourney v6, GPT-4o) are more photorealistic.",
        "Critical question: do existing detectors still work — and why/why not?",
    ])

    add_bullet_slide(prs, "Research Question & Hypothesis", [
        "RQ: How well do benchmark-trained detectors generalise to next-gen generators?",
        "Hypothesis: models rely on generator-specific low-level artefacts.",
        "Those artefacts weaken in newer models → measurable generalisation gap.",
    ])

    add_bullet_slide(prs, "Related Work & Gap", [
        "Wang 2020: GAN→GAN transfer with augmentation.",
        "Corvi 2023: GAN↔diffusion gap + frequency analysis.",
        "Bird & Lotfi 2024: CIFAKE benchmark.",
        "Our gap: multi-family eval + calibration + Grad-CAM + measure + resolution control.",
    ])

    add_bullet_slide(prs, "Approach", [
        "Fine-tune EfficientNet-B3 on CIFAKE (SD v1.4).",
        "Temperature scaling (T = 1.2189) + Grad-CAM (conv_head).",
        "Evaluate StyleGAN, SD3/Flux, Midjourney v6, GPT-4o, Janus-Pro (300 each).",
        "Diagnose GPT-4o with FFT, t-SNE, quantitative Grad-CAM.",
        "Resolution-matched real-image control (notebook 07).",
    ])

    add_bullet_slide(
        prs, "Baseline Results (CIFAKE)",
        [
            "Accuracy 96.96%  ·  AUC 0.9971  ·  ECE 0.0026 (pre-T).",
            "FAKE recall 95.3% / REAL recall 98.6%.",
            "Temperature scaling slightly raised ECE (NLL≠ECE).",
        ],
        PLOTS / "roc_curve.png",
    )

    add_bullet_slide(
        prs, "Calibration",
        [
            "Learned T = 1.2189 on validation set (L-BFGS).",
            "Already well-calibrated in-distribution.",
            "Under OOD generators, CIFAKE-fit T worsens ECE.",
        ],
        PLOTS / "calibration_comparison.png",
    )

    add_bullet_slide(
        prs, "Grad-CAM",
        [
            "REAL: diffuse texture attention, high confidence.",
            "FAKE: localised artefact hotspots.",
            "Suggests SD v1.4-specific cues — motivates generalisation study.",
        ],
        PLOTS / "gradcam_comparison.png",
    )

    add_bullet_slide(prs, "Cross-Generator Design", [
        "Five families × 300 fakes + 300 CIFAKE reals.",
        "Primary OOD metric: fake detection rate (not overall accuracy).",
        "Caveat: resolution confound → Experiment 5 control.",
    ])

    add_bullet_slide(
        prs, "Cross-Generator Results",
        [
            "StyleGAN / MJ / Janus-Pro: ~94% fake detection.",
            "SD3/Flux: 97% fake detection (strong transfer).",
            "GPT-4o: 86.3% fake detection — only meaningful gap (4.8% acc. drop).",
        ],
        PLOTS / "cross_generator_accuracy.png",
    )

    add_bullet_slide(
        prs, "Hypothesis Revised",
        [
            "Naive protocol: transfer looks strong except GPT-4o.",
            "Control: 93% of high-res reals called FAKE → confound.",
            "After matching: broad gap (fake detection ~35–62%).",
            "GPT-4o feature-absence remains a secondary native-protocol finding.",
        ],
        PLOTS / "degradation_waterfall.png",
    )

    gpt_dir = PLOTS / "gpt4o_investigation"
    add_bullet_slide(
        prs, "Why GPT-4o Fails",
        [
            "FFT: spectra closer to real photos.",
            "t-SNE: embeddings near real cluster.",
            "Grad-CAM metrics: not uniquely unfocused vs Midjourney/SD3.",
            "Conclusion: feature absence, not attention failure.",
        ],
        gpt_dir / "feature_space_tsne.png",
    )

    ctrl = PLOTS / "resolution_control" / "fake_rate_by_condition.png"
    add_bullet_slide(
        prs, "Resolution Control — Confound Confirmed",
        [
            "A CIFAKE REAL: 1.7% predicted FAKE.",
            "B Hi-res REAL: 93.0% predicted FAKE — not generator detection.",
            "C Same photos →32×32: FAKE-rate falls to 37.3%.",
            "Matched fakes collapse: Midjourney 94%→35%, StyleGAN 94%→38%.",
            "Experiment 3 transfer was largely resolution-driven.",
        ],
        ctrl if ctrl.exists() else None,
    )

    add_bullet_slide(prs, "Live Demo", [
        "Gradio app: classification + calibrated confidence + Grad-CAM.",
        "Samples: CIFAKE real/fake, StyleGAN, Midjourney, GPT-4o.",
        "python app.py   or   python app.py --share",
        "See reports/demo_script.md for the spoken runbook.",
    ])

    add_bullet_slide(prs, "Limitations & Future Work", [
        "Single-generator training; 32×32 CIFAKE origin.",
        "N=300/family; single backbone.",
        "Future: multi-generator training, higher-res data, frequency features, recalibration.",
    ])

    add_bullet_slide(prs, "Conclusions", [
        "Strong CIFAKE detector in-distribution (96.96%, AUC 0.9971).",
        "Naive cross-generator transfer was resolution-confounded.",
        "Matched-resolution eval reveals a broad generalisation gap.",
        "GPT-4o: feature absence under native protocol (FFT / t-SNE / Grad-CAM).",
        "Low-res benchmarks need matched-resolution real controls.",
        "Code, tests (71), report, demo: github.com/krishi-shah/ai-image-detection",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
